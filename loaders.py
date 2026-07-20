from __future__ import annotations

import csv
import json
import os
import re
from pathlib import Path
from typing import Iterable

import docx
import markitdown
import openpyxl
import pymupdf
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".csv", ".md", ".txt", ".json", ".html"}
DEFAULT_CHUNK_SIZE = 450
DEFAULT_CHUNK_OVERLAP = 60


def _normalize_section_name(name: str) -> str:
    name = re.sub(r"^\s*#+\s*", "", name).strip()
    name = re.sub(r"^\s*[-*]\s*", "", name).strip()
    return name or "未命名章节"


def _extract_sections(text: str) -> tuple[list[dict], str]:
    sections: list[dict] = []
    current_title = "文档正文"
    current_start = 0
    lines = text.splitlines()
    for offset, line in enumerate(lines):
        if re.match(r"^(#{1,6}\s+.+|[-*]\s+.+)$", line):
            sections.append({"title": current_title, "start": current_start, "end": sum(len(l) + 1 for l in lines[:offset])})
            current_title = _normalize_section_name(line)
            current_start = sum(len(l) + 1 for l in lines[:offset])
    sections.append({"title": current_title, "start": current_start, "end": len(text)})
    return sections, text


def _assign_section(char_index: int, sections: list[dict]) -> str:
    for section in reversed(sections):
        if char_index >= section["start"]:
            return section["title"]
    return sections[-1]["title"] if sections else "文档正文"


def _split_with_positions(text: str, chunk_size: int = DEFAULT_CHUNK_SIZE, chunk_overlap: int = DEFAULT_CHUNK_OVERLAP) -> Iterable[dict]:
    sections, _ = _extract_sections(text)
    cleaned = re.sub(r"\s+", " ", text).strip()
    if not cleaned:
        return
    if len(cleaned) <= chunk_size:
        yield {
            "text": cleaned,
            "chunk_id": "chunk_001",
            "char_start": 0,
            "char_end": len(cleaned),
            "section": _assign_section(0, sections),
        }
        return
    chunk_index = 0
    start = 0
    while start < len(cleaned):
        end = min(start + chunk_size, len(cleaned))
        chunk_text = cleaned[start:end]
        chunk_index += 1
        yield {
            "text": chunk_text,
            "chunk_id": f"chunk_{chunk_index:03d}",
            "char_start": start,
            "char_end": end,
            "section": _assign_section(start, sections),
        }
        if end == len(cleaned):
            break
        start = end - chunk_overlap
        if start < 0:
            start = 0


def build_documents_from_file(path: Path) -> list[dict]:
    documents = []
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        documents.extend(load_pdf(path))
    elif suffix in {".docx", ".doc"}:
        documents.extend(load_docx(path))
    elif suffix in {".pptx", ".ppt"}:
        documents.extend(load_pptx(path))
    elif suffix in {".xlsx", ".xls"}:
        documents.extend(load_xlsx(path))
    elif suffix == ".csv":
        documents.extend(load_csv(path))
    elif suffix == ".json":
        documents.extend(load_json_file(path))
    elif suffix in {".md", ".txt", ".html"}:
        documents.extend(load_text_like(path))
    else:
        documents.extend(load_markitdown_fallback(path))
    return documents


def load_pdf(path: Path) -> Iterable[dict]:
    try:
        document = pymupdf.open(path)
        for page_index, page in enumerate(document, start=1):
            text = page.get_text("text")
            if not text.strip():
                continue
            prefix = f"[Page {page_index}] "
            for chunk in _split_with_positions(prefix + text):
                chunk["source"] = path.name
                chunk["page"] = page_index
                yield chunk
    except Exception:  # noqa: BLE001
        yield from load_markitdown_fallback(path)


def load_docx(path: Path) -> Iterable[dict]:
    try:
        document = docx.Document(path)
        text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)
        if text.strip():
            table_texts = []
            for table in document.tables:
                for row in table.rows:
                    table_texts.append(" | ".join(cell.text for cell in row.cells))
            if table_texts:
                text += "\n\n" + "\n".join(table_texts)
        if not text.strip():
            return
        for chunk in _split_with_positions(text):
            chunk["source"] = path.name
            chunk["page"] = 1
            yield chunk
    except Exception:  # noqa: BLE001
        yield from load_markitdown_fallback(path)


def load_pptx(path: Path) -> Iterable[dict]:
    try:
        presentation = Presentation(path)
        for page_index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            text = "\n".join(texts)
            if not text:
                continue
            prefix = f"[Page {page_index}] "
            for chunk in _split_with_positions(prefix + text):
                chunk["source"] = path.name
                chunk["page"] = page_index
                yield chunk
    except Exception:  # noqa: BLE001
        yield from load_markitdown_fallback(path)


def load_xlsx(path: Path) -> Iterable[dict]:
    try:
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
        chunks = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                if any(value is not None for value in row):
                    rows.append(" | ".join("" if value is None else str(value) for value in row))
            if rows:
                chunks.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        workbook.close()
        text = "\n\n".join(chunks)
        if not text.strip():
            return
        for chunk in _split_with_positions(text):
            chunk["source"] = path.name
            chunk["page"] = 1
            yield chunk
    except Exception:  # noqa: BLE001
        yield from load_markitdown_fallback(path)


def load_csv(path: Path) -> Iterable[dict]:
    try:
        with path.open("r", encoding="utf-8", errors="replace") as file:
            reader = csv.reader(file)
            rows = [" | ".join(row) for row in reader]
        text = "\n".join(rows)
        if not text.strip():
            return
        for chunk in _split_with_positions(text):
            chunk["source"] = path.name
            chunk["page"] = 1
            yield chunk
    except Exception:  # noqa: BLE001
        yield from load_markitdown_fallback(path)


def load_json_file(path: Path) -> Iterable[dict]:
    try:
        with path.open("r", encoding="utf-8", errors="replace") as file:
            data = json.load(file)
        text = json.dumps(data, ensure_ascii=False, indent=2)
        for chunk in _split_with_positions(text):
            chunk["source"] = path.name
            chunk["page"] = 1
            yield chunk
    except Exception:  # noqa: BLE001
        yield from load_markitdown_fallback(path)


def load_text_like(path: Path) -> Iterable[dict]:
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except Exception:  # noqa: BLE001
        return
    for chunk in _split_with_positions(text):
        chunk["source"] = path.name
        chunk["page"] = 1
        yield chunk


def load_markitdown_fallback(path: Path) -> Iterable[dict]:
    try:
        converter = markitdown.MarkItDown(enable_plugins=False)
        result = converter.convert(str(path))
        if result and getattr(result, "text_content", None):
            for chunk in _split_with_positions(result.text_content):
                chunk["source"] = path.name
                chunk["page"] = 1
                yield chunk
    except Exception:  # noqa: BLE001
        return


def split_text(text: str, chunk_size: int = DEFAULT_CHUNK_SIZE, chunk_overlap: int = DEFAULT_CHUNK_OVERLAP) -> list[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    if len(text) <= chunk_size:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - chunk_overlap
        if start < 0:
            start = 0
    return [chunk for chunk in chunks if chunk.strip()]
